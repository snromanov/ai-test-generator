#!/usr/bin/env python3
"""
AI Test Generator - Executive Presentation Generator
Creates a professional PowerPoint presentation for bank leadership.
"""

import json
import os
from pathlib import Path
from datetime import datetime

# Charts generation
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Colors (Corporate Bank Style)
COLORS = {
    'primary_blue': RGBColor(0x00, 0x33, 0x66),      # #003366 - Dark Blue
    'secondary_blue': RGBColor(0x00, 0x66, 0xCC),    # #0066CC - Medium Blue
    'light_blue': RGBColor(0xE6, 0xF0, 0xFA),        # #E6F0FA - Light Blue
    'gray': RGBColor(0x6C, 0x75, 0x7D),              # #6C757D - Gray
    'dark_gray': RGBColor(0x34, 0x34, 0x34),         # #343434 - Dark Gray
    'white': RGBColor(0xFF, 0xFF, 0xFF),             # #FFFFFF - White
    'red': RGBColor(0xDC, 0x35, 0x45),               # #DC3545 - Red (Risk)
    'yellow': RGBColor(0xFF, 0xC1, 0x07),            # #FFC107 - Yellow (Warning)
    'green': RGBColor(0x28, 0xA7, 0x45),             # #28A745 - Green (Success)
    'orange': RGBColor(0xFD, 0x7E, 0x14),            # #FD7E14 - Orange
}

# Matplotlib colors
MPL_COLORS = {
    'primary_blue': '#003366',
    'secondary_blue': '#0066CC',
    'light_blue': '#E6F0FA',
    'gray': '#6C757D',
    'red': '#DC3545',
    'green': '#28A745',
    'yellow': '#FFC107',
    'orange': '#FD7E14',
}

# Paths
BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / 'data'
CHARTS_DIR = BASE_DIR / 'charts'
OUTPUT_DIR = BASE_DIR / 'output'


def load_data():
    """Load JSON data files."""
    with open(DATA_DIR / 'downtime_costs.json', 'r', encoding='utf-8') as f:
        downtime_data = json.load(f)
    with open(DATA_DIR / 'roi_calculations.json', 'r', encoding='utf-8') as f:
        roi_data = json.load(f)
    return downtime_data, roi_data


def create_charts(downtime_data, roi_data):
    """Generate chart images for the presentation."""
    CHARTS_DIR.mkdir(exist_ok=True)

    # Set font for Russian text
    plt.rcParams['font.family'] = ['DejaVu Sans', 'Arial', 'sans-serif']
    plt.rcParams['font.size'] = 12

    # Chart 1: Time Comparison (Manual vs AI)
    fig, ax = plt.subplots(figsize=(10, 6))
    categories = ['Ручная работа', 'AI Generator']
    times = [17 * 60, 5]  # Convert to minutes for comparison
    colors = [MPL_COLORS['gray'], MPL_COLORS['secondary_blue']]

    bars = ax.barh(categories, times, color=colors, height=0.5)
    ax.set_xlabel('Время (минуты)', fontsize=14, fontweight='bold')
    ax.set_title('Генерация тест-кейсов: Manual vs AI', fontsize=16, fontweight='bold', pad=20)

    # Add value labels
    for bar, time in zip(bars, [17*60, 5]):
        width = bar.get_width()
        label = f'{17} часов' if time > 60 else f'{time} мин'
        ax.text(width + 20, bar.get_y() + bar.get_height()/2, label,
                va='center', fontsize=14, fontweight='bold')

    ax.set_xlim(0, 1200)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # Add efficiency label
    ax.text(600, -0.5, 'В 204 раза быстрее', fontsize=18, fontweight='bold',
            color=MPL_COLORS['green'], ha='center')

    plt.tight_layout()
    plt.savefig(CHARTS_DIR / 'time_comparison.png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    # Chart 2: ROI Comparison (3 years)
    fig, ax = plt.subplots(figsize=(10, 6))

    years = ['Год 1', 'Год 2', 'Год 3']
    hire_qa_costs = [4.98, 5.15, 5.66]  # In millions RUB
    ai_costs = [0.41, 0.37, 0.41]       # In millions RUB

    x = np.arange(len(years))
    width = 0.35

    bars1 = ax.bar(x - width/2, hire_qa_costs, width, label='Найм QA', color=MPL_COLORS['red'])
    bars2 = ax.bar(x + width/2, ai_costs, width, label='AI Generator', color=MPL_COLORS['green'])

    ax.set_ylabel('Стоимость (млн руб)', fontsize=14, fontweight='bold')
    ax.set_title('Сравнение затрат: Найм QA vs AI Generator', fontsize=16, fontweight='bold', pad=20)
    ax.set_xticks(x)
    ax.set_xticklabels(years, fontsize=12)
    ax.legend(fontsize=12)

    # Add value labels on bars
    for bar in bars1:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:.2f}M', ha='center', va='bottom', fontsize=11, fontweight='bold')
    for bar in bars2:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:.2f}M', ha='center', va='bottom', fontsize=11, fontweight='bold')

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_ylim(0, 7)

    plt.tight_layout()
    plt.savefig(CHARTS_DIR / 'roi_comparison.png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    # Chart 3: Downtime Cost Growth
    fig, ax = plt.subplots(figsize=(10, 6))

    hours = [1, 2, 4, 8, 24]
    cost_min = [0.3, 0.6, 1.2, 2.4, 7.2]   # $300k/hour minimum
    cost_max = [5.0, 10.0, 20.0, 40.0, 120.0]  # $5M/hour for finance

    ax.fill_between(hours, cost_min, cost_max, alpha=0.3, color=MPL_COLORS['red'])
    ax.plot(hours, cost_min, 'o-', color=MPL_COLORS['orange'], linewidth=2, label='Минимум ($300k/час)')
    ax.plot(hours, cost_max, 'o-', color=MPL_COLORS['red'], linewidth=2, label='Финансовый сектор ($5M/час)')

    ax.set_xlabel('Длительность простоя (часы)', fontsize=14, fontweight='bold')
    ax.set_ylabel('Стоимость простоя (млн USD)', fontsize=14, fontweight='bold')
    ax.set_title('Рост стоимости простоя со временем', fontsize=16, fontweight='bold', pad=20)
    ax.legend(fontsize=11, loc='upper left')

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_xlim(0, 25)
    ax.set_ylim(0, 130)

    # Add annotation
    ax.annotate('Knight Capital:\n45 мин = $440M', xy=(0.75, 440/10), xytext=(5, 80),
                fontsize=12, fontweight='bold',
                arrowprops=dict(arrowstyle='->', color=MPL_COLORS['red']),
                color=MPL_COLORS['red'])

    plt.tight_layout()
    plt.savefig(CHARTS_DIR / 'downtime_cost_growth.png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    # Chart 4: Risk Timeline
    fig, ax = plt.subplots(figsize=(12, 4))

    timeline = ['Сегодня', 'Месяц 1\n(пилот)', 'Месяц 3', 'Месяц 6']
    risk_levels = [100, 60, 30, 10]  # Risk percentage
    colors_timeline = [MPL_COLORS['red'], MPL_COLORS['yellow'], MPL_COLORS['green'], MPL_COLORS['green']]

    for i, (label, risk, color) in enumerate(zip(timeline, risk_levels, colors_timeline)):
        circle = plt.Circle((i * 2, 0), 0.4, color=color, zorder=3)
        ax.add_patch(circle)
        ax.text(i * 2, 0, f'{risk}%', ha='center', va='center', fontsize=14,
                fontweight='bold', color='white', zorder=4)
        ax.text(i * 2, -0.9, label, ha='center', va='top', fontsize=12, fontweight='bold')

        if i < len(timeline) - 1:
            ax.arrow(i * 2 + 0.5, 0, 1, 0, head_width=0.15, head_length=0.1,
                    fc=MPL_COLORS['gray'], ec=MPL_COLORS['gray'], zorder=2)

    ax.set_xlim(-1, 7)
    ax.set_ylim(-1.5, 1)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.set_title('Снижение риска с течением времени', fontsize=16, fontweight='bold', pad=20)

    plt.tight_layout()
    plt.savefig(CHARTS_DIR / 'risk_timeline.png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()

    print("Charts created successfully!")


def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size=44,
                    font_color=COLORS['primary_blue'], bold=True, alignment=PP_ALIGN.LEFT):
    """Add a title text box to slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return shape


def add_body_text(slide, text, left, top, width, height, font_size=18,
                  font_color=COLORS['dark_gray'], bold=False):
    """Add body text box to slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    return shape


def add_bullet_points(slide, items, left, top, width, height, font_size=16,
                      font_color=COLORS['dark_gray']):
    """Add bulleted list to slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = Pt(8)
    return shape


def add_big_number(slide, number, label, left, top, width, height,
                   number_color=COLORS['primary_blue'], font_size=72):
    """Add a large number with label below."""
    # Number
    shape = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(font_size)
    p.font.color.rgb = number_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_shape = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(0.5))
    tf = label_shape.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return shape, label_shape


def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rectangle(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def create_slide_1_title(prs, downtime_data, roi_data):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient effect (solid color for simplicity)
    set_slide_background(slide, COLORS['primary_blue'])

    # Main title
    add_title_shape(slide, "AI Test Generator",
                    Inches(0.5), Inches(2), Inches(12), Inches(1),
                    font_size=54, font_color=COLORS['white'], bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Subtitle line 1
    add_title_shape(slide, "Решение критического риска в Infrastructure",
                    Inches(0.5), Inches(2.9), Inches(12), Inches(0.6),
                    font_size=28, font_color=COLORS['light_blue'], bold=False,
                    alignment=PP_ALIGN.CENTER)

    # Key value proposition
    add_rounded_rectangle(slide, Inches(2.5), Inches(4), Inches(8), Inches(0.8), COLORS['yellow'])
    add_title_shape(slide, "4 недели до покрытия vs 9 месяцев найма QA",
                    Inches(2.5), Inches(4.15), Inches(8), Inches(0.6),
                    font_size=24, font_color=COLORS['dark_gray'], bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Date
    add_body_text(slide, "Январь 2025",
                  Inches(0.5), Inches(6.8), Inches(4), Inches(0.4),
                  font_size=16, font_color=COLORS['light_blue'])

    # Placeholder for logo area
    add_body_text(slide, "Райффайзенбанк",
                  Inches(9), Inches(6.8), Inches(3.5), Inches(0.4),
                  font_size=16, font_color=COLORS['light_blue'])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Добро пожаловать. Сегодня представлю решение критического риска в инфраструктуре.
- Ключевое сообщение: мы можем получить тестовое покрытие за 4 недели вместо 9 месяцев найма.
- Это не теоретическое решение - MVP уже готов и протестирован.
"""


def create_slide_2_problem(prs, downtime_data, roi_data):
    """Slide 2: Problem Statement"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Критический риск обнаружен",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=40, font_color=COLORS['red'])

    # Warning banner
    warning_box = add_rounded_rectangle(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2), COLORS['red'])
    add_title_shape(slide, "Инфраструктура БЕЗ QA = Фундамент из песка",
                    Inches(0.7), Inches(1.5), Inches(11.5), Inches(0.8),
                    font_size=32, font_color=COLORS['white'], bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Key facts grid
    facts = [
        ("Направление", "Infrastructure\n(CI/CD, Monitoring, API Gateway)"),
        ("Текущее покрытие", "0%"),
        ("Планы по найму QA", "Нет"),
        ("Зависимости", "ВСЕ продукты банка"),
    ]

    for i, (label, value) in enumerate(facts):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.5)
        top = Inches(2.8 + row * 1.8)

        # Box
        add_rounded_rectangle(slide, left, top, Inches(6), Inches(1.5), COLORS['light_blue'])

        # Label
        add_body_text(slide, label, left + Inches(0.2), top + Inches(0.15),
                      Inches(5.6), Inches(0.4), font_size=14, font_color=COLORS['gray'])

        # Value
        color = COLORS['red'] if value in ['0%', 'Нет'] else COLORS['primary_blue']
        add_title_shape(slide, value, left + Inches(0.2), top + Inches(0.5),
                        Inches(5.6), Inches(0.9), font_size=24, font_color=color)

    # Bottom warning
    add_body_text(slide, "Любой баг в инфраструктуре = потенциальный простой всего банка",
                  Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                  font_size=18, font_color=COLORS['red'], bold=True)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Инфраструктура - это фундамент, на котором работают все продукты банка.
- Сейчас это направление работает полностью без QA покрытия.
- Это означает, что мы полагаемся только на happy path тестирование самими разработчиками.
- При этом любой баг в CI/CD или API Gateway может привести к каскадному отказу всех систем.
"""


def create_slide_3_downtime_cost(prs, downtime_data, roi_data):
    """Slide 3: Cost of Downtime"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Стоимость инфраструктурных инцидентов",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=36, font_color=COLORS['primary_blue'])

    # Big numbers row
    numbers_data = [
        ("$5,600", "в минуту", "средний простой\n(Gartner)"),
        ("$5M+", "в час", "финансовый сектор\n(Gartner 2024)"),
        ("$152M", "в год", "потери от простоев\n(Splunk/Oxford)"),
    ]

    for i, (number, unit, source) in enumerate(numbers_data):
        left = Inches(0.5 + i * 4.3)
        top = Inches(1.2)

        # Box
        color = COLORS['red'] if i == 1 else COLORS['light_blue']
        add_rounded_rectangle(slide, left, top, Inches(4), Inches(2.2), color)

        # Number
        text_color = COLORS['white'] if i == 1 else COLORS['primary_blue']
        add_title_shape(slide, number, left + Inches(0.1), top + Inches(0.2),
                        Inches(3.8), Inches(0.9), font_size=48, font_color=text_color,
                        alignment=PP_ALIGN.CENTER)

        # Unit
        add_body_text(slide, unit, left + Inches(0.1), top + Inches(1.1),
                      Inches(3.8), Inches(0.4), font_size=20,
                      font_color=COLORS['white'] if i == 1 else COLORS['gray'], bold=True)
        shape = slide.shapes[-1]
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Source
        add_body_text(slide, source, left + Inches(0.1), top + Inches(1.6),
                      Inches(3.8), Inches(0.5), font_size=12,
                      font_color=COLORS['white'] if i == 1 else COLORS['gray'])
        shape = slide.shapes[-1]
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Chart
    chart_path = CHARTS_DIR / 'downtime_cost_growth.png'
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(3.6), width=Inches(8))

    # Sources box
    add_rounded_rectangle(slide, Inches(8.7), Inches(3.6), Inches(4.1), Inches(3.2), COLORS['light_blue'])
    sources_text = """Источники данных:

• Gartner Research 2024
• Uptime Institute Annual Report
• ITIC 2024 Survey
• Splunk/Oxford Economics

90% компаний теряют
>$300,000 за час простоя"""

    add_body_text(slide, sources_text, Inches(8.9), Inches(3.8),
                  Inches(3.7), Inches(2.8), font_size=13, font_color=COLORS['dark_gray'])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Это не наши внутренние оценки - это данные от ведущих аналитических агентств.
- Gartner: средняя стоимость простоя $5,600 в минуту, для финансового сектора - более $5M в час.
- По данным Splunk и Oxford Economics, финансовые организации теряют $152M в год из-за простоев.
- 90% компаний по данным ITIC теряют более $300,000 за каждый час простоя.
- Обратите внимание на график - стоимость растёт экспоненциально со временем.
"""


def create_slide_4_real_incidents(prs, downtime_data, roi_data):
    """Slide 4: Real Incident Scenarios"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Реальные инциденты: что может пойти не так",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=36, font_color=COLORS['primary_blue'])

    incidents = [
        {
            "company": "Knight Capital (2012)",
            "cause": "Баг в деплое: 7 из 8 серверов обновлены",
            "duration": "45 минут",
            "loss": "$440M",
            "result": "Компания продана через 4 месяца",
            "color": COLORS['red']
        },
        {
            "company": "TSB Bank (2018)",
            "cause": "Провал миграции IT-систем",
            "duration": "8 месяцев",
            "loss": "£330M",
            "result": "CEO уволен, £49M штрафов",
            "color": COLORS['orange']
        },
        {
            "company": "RBS/NatWest (2012)",
            "cause": "Откат апгрейда без тестирования",
            "duration": "3 недели",
            "loss": "£231M",
            "result": "£56M штрафов, 6.5M клиентов",
            "color": COLORS['yellow']
        }
    ]

    for i, incident in enumerate(incidents):
        top = Inches(1.3 + i * 1.8)

        # Left color bar
        add_rectangle(slide, Inches(0.5), top, Inches(0.15), Inches(1.6), incident['color'])

        # Content box
        add_rounded_rectangle(slide, Inches(0.7), top, Inches(12.1), Inches(1.6), COLORS['light_blue'])

        # Company name
        add_title_shape(slide, incident['company'], Inches(0.9), top + Inches(0.1),
                        Inches(4), Inches(0.5), font_size=20, font_color=COLORS['primary_blue'])

        # Cause
        add_body_text(slide, f"Причина: {incident['cause']}", Inches(0.9), top + Inches(0.55),
                      Inches(5), Inches(0.4), font_size=14, font_color=COLORS['dark_gray'])

        # Duration
        add_body_text(slide, f"Длительность: {incident['duration']}", Inches(0.9), top + Inches(0.95),
                      Inches(3), Inches(0.4), font_size=14, font_color=COLORS['gray'])

        # Result
        add_body_text(slide, incident['result'], Inches(0.9), top + Inches(1.25),
                      Inches(5), Inches(0.4), font_size=12, font_color=COLORS['gray'])

        # Loss - big number
        add_title_shape(slide, incident['loss'], Inches(9), top + Inches(0.3),
                        Inches(3.5), Inches(0.8), font_size=40, font_color=incident['color'],
                        alignment=PP_ALIGN.RIGHT)
        add_body_text(slide, "потери", Inches(9), top + Inches(1.1),
                      Inches(3.5), Inches(0.3), font_size=14, font_color=COLORS['gray'])
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Bottom insight
    add_rounded_rectangle(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6), COLORS['primary_blue'])
    add_body_text(slide, "Все три инцидента могли быть предотвращены качественным тестированием инфраструктуры",
                  Inches(0.7), Inches(6.42), Inches(12), Inches(0.4),
                  font_size=16, font_color=COLORS['white'], bold=True)
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Knight Capital: баг в деплое - только 7 из 8 серверов получили обновление. За 45 минут потеряли $440M.
- TSB Bank: миграция на новую платформу без достаточного тестирования. 1.9M клиентов не могли войти в аккаунты.
- RBS/NatWest: откатили обновление без тестирования совместимости. 6.5 миллионов клиентов пострадали.
- Общий знаменатель: недостаточное тестирование инфраструктурных изменений.
- Это именно тот риск, который мы сейчас несём в нашем Infrastructure направлении.
"""


def create_slide_5_traditional_solutions(prs, downtime_data, roi_data):
    """Slide 5: Why Traditional Solutions Don't Work"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Почему традиционные подходы не работают",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=36, font_color=COLORS['primary_blue'])

    # Table header
    headers = ["Решение", "Timeline", "Стоимость", "Проблема"]
    header_widths = [3.5, 2, 2, 5.3]
    header_left = 0.5

    # Header row
    for i, (header, width) in enumerate(zip(headers, header_widths)):
        add_rectangle(slide, Inches(header_left), Inches(1.3), Inches(width), Inches(0.6), COLORS['primary_blue'])
        add_body_text(slide, header, Inches(header_left + 0.1), Inches(1.4),
                      Inches(width - 0.2), Inches(0.4), font_size=14, font_color=COLORS['white'], bold=True)
        header_left += width

    # Table data
    rows = [
        ["Нанять Infrastructure QA", "6-9 месяцев", "~5M₽/год", "Редкий специалист, долгий поиск"],
        ["Перераспределить QA", "1 месяц", "0₽", "Создаст новые дыры в покрытии"],
        ["Инженеры тестируют сами", "0", "0₽", "Только happy path, bias автора"],
        ["Аутсорс тестирования", "2-3 месяца", "3-4M₽/год", "Нет знания специфики банка, security риски"],
    ]

    for row_idx, row in enumerate(rows):
        header_left = 0.5
        row_top = Inches(1.95 + row_idx * 0.9)
        bg_color = COLORS['light_blue'] if row_idx % 2 == 0 else COLORS['white']

        for col_idx, (cell, width) in enumerate(zip(row, header_widths)):
            add_rectangle(slide, Inches(header_left), row_top, Inches(width), Inches(0.85), bg_color)

            font_color = COLORS['red'] if col_idx == 3 else COLORS['dark_gray']
            add_body_text(slide, cell, Inches(header_left + 0.15), row_top + Inches(0.25),
                          Inches(width - 0.3), Inches(0.5), font_size=13, font_color=font_color)
            header_left += width

    # Conclusion box
    add_rounded_rectangle(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), COLORS['yellow'])

    add_title_shape(slide, "Вывод:", Inches(0.7), Inches(5.75),
                    Inches(2), Inches(0.5), font_size=20, font_color=COLORS['dark_gray'])

    add_body_text(slide, """• Традиционные решения либо слишком долгие (6-9 месяцев), либо создают новые риски
• Каждый месяц без покрытия = накопление технического долга
• Нужно решение, которое работает СЕЙЧАС""",
                  Inches(0.7), Inches(6.15), Inches(12), Inches(0.6),
                  font_size=14, font_color=COLORS['dark_gray'])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Нанять Infrastructure QA: рынок пустой, специалисты редкие, 6-9 месяцев на поиск + онбординг.
- Перераспределить существующих QA: создаём новые дыры в покрытии других направлений.
- Инженеры тестируют сами: классическая проблема - автор кода не видит своих ошибок, тестирует только happy path.
- Аутсорс: нет понимания специфики банковских систем, плюс security риски с доступом.
- Нам нужно решение, которое может работать уже сейчас, без долгого ожидания.
"""


def create_slide_6_solution(prs, downtime_data, roi_data):
    """Slide 6: Our Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "AI Test Generator",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=44, font_color=COLORS['primary_blue'])

    add_body_text(slide, "Автоматическая генерация тест-кейсов из требований",
                  Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
                  font_size=18, font_color=COLORS['gray'])

    # Workflow visualization
    workflow_steps = [
        ("Confluence\nRequirements", COLORS['gray']),
        ("AI\nAnalysis", COLORS['secondary_blue']),
        ("Test\nCases", COLORS['green']),
        ("Engineers\nExecute", COLORS['primary_blue']),
    ]

    for i, (step, color) in enumerate(workflow_steps):
        left = Inches(0.8 + i * 3.2)

        # Circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(1.5), Inches(1.8), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Text in circle
        add_body_text(slide, step, left + Inches(0.1), Inches(2.1),
                      Inches(1.6), Inches(0.8), font_size=14, font_color=COLORS['white'], bold=True)
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < len(workflow_steps) - 1:
            arrow_left = left + Inches(1.85)
            add_body_text(slide, "→", arrow_left, Inches(2.1),
                          Inches(0.5), Inches(0.5), font_size=36, font_color=COLORS['gray'], bold=True)

    # Features grid
    features = [
        ("5 минут", "генерации vs 17 часов\nручной работы", "time"),
        ("8-12", "тест-кейсов на\nтребование", "coverage"),
        ("On-premise", "данные не покидают\nбанк", "security"),
        ("MVP готов", "можем начать\nпилот сегодня", "ready"),
    ]

    for i, (number, desc, icon) in enumerate(features):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.5)
        top = Inches(3.7 + row * 1.5)

        # Box
        add_rounded_rectangle(slide, left, top, Inches(6), Inches(1.3), COLORS['light_blue'])

        # Icon placeholder (using emoji/symbol)
        icons = {"time": "⚡", "coverage": "🎯", "security": "🔒", "ready": "🚀"}
        add_title_shape(slide, icons[icon], left + Inches(0.2), top + Inches(0.3),
                        Inches(0.8), Inches(0.8), font_size=32, font_color=COLORS['primary_blue'])

        # Number/Title
        add_title_shape(slide, number, left + Inches(1), top + Inches(0.15),
                        Inches(2.5), Inches(0.6), font_size=28, font_color=COLORS['primary_blue'])

        # Description
        add_body_text(slide, desc, left + Inches(1), top + Inches(0.7),
                      Inches(4.5), Inches(0.6), font_size=14, font_color=COLORS['gray'])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Решение: AI Test Generator - автоматическая генерация тест-кейсов из требований.
- Workflow простой: берём требования из Confluence, AI анализирует и генерирует тест-кейсы, инженеры выполняют.
- Ключевые преимущества:
  - Скорость: 5 минут вместо 17 часов ручной работы
  - Полнота: 8-12 тест-кейсов на каждое требование, включая edge cases
  - Безопасность: on-premise развёртывание, данные не покидают банк
  - Готовность: MVP уже готов, можем начать пилот сегодня
"""


def create_slide_7_demo_results(prs, downtime_data, roi_data):
    """Slide 7: Demo Results"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Proof of Concept: PetStore API",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=36, font_color=COLORS['primary_blue'])

    # Input section
    add_rounded_rectangle(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2), COLORS['light_blue'])
    add_title_shape(slide, "INPUT", Inches(0.7), Inches(1.35),
                    Inches(2), Inches(0.4), font_size=18, font_color=COLORS['gray'])

    add_title_shape(slide, "9", Inches(0.7), Inches(1.8),
                    Inches(2), Inches(0.7), font_size=48, font_color=COLORS['primary_blue'])
    add_body_text(slide, "API endpoints", Inches(0.7), Inches(2.5),
                  Inches(2), Inches(0.3), font_size=14, font_color=COLORS['gray'])

    add_title_shape(slide, "2", Inches(3), Inches(1.8),
                    Inches(2), Inches(0.7), font_size=48, font_color=COLORS['primary_blue'])
    add_body_text(slide, "страницы документации", Inches(3), Inches(2.5),
                  Inches(2.5), Inches(0.3), font_size=14, font_color=COLORS['gray'])

    # Arrow
    add_title_shape(slide, "→", Inches(5.8), Inches(1.9),
                    Inches(1), Inches(0.8), font_size=48, font_color=COLORS['green'],
                    alignment=PP_ALIGN.CENTER)

    # Output section
    add_rounded_rectangle(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(2), COLORS['green'])
    add_title_shape(slide, "OUTPUT", Inches(7.2), Inches(1.35),
                    Inches(2), Inches(0.4), font_size=18, font_color=COLORS['white'])

    add_title_shape(slide, "52", Inches(7.2), Inches(1.8),
                    Inches(3), Inches(0.7), font_size=56, font_color=COLORS['white'])
    add_body_text(slide, "тест-кейса", Inches(7.2), Inches(2.5),
                  Inches(2), Inches(0.3), font_size=16, font_color=COLORS['white'], bold=True)

    add_title_shape(slide, "5 мин", Inches(10), Inches(1.8),
                    Inches(2.3), Inches(0.7), font_size=32, font_color=COLORS['white'])
    add_body_text(slide, "время генерации", Inches(10), Inches(2.5),
                  Inches(2.3), Inches(0.3), font_size=12, font_color=COLORS['white'])

    # Chart
    chart_path = CHARTS_DIR / 'time_comparison.png'
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(3.5), width=Inches(7.5))

    # Coverage breakdown
    add_rounded_rectangle(slide, Inches(8.3), Inches(3.5), Inches(4.2), Inches(3.3), COLORS['light_blue'])
    add_title_shape(slide, "Покрытие:", Inches(8.5), Inches(3.65),
                    Inches(3.5), Inches(0.4), font_size=18, font_color=COLORS['primary_blue'])

    coverage_items = [
        "✓ Позитивные сценарии",
        "✓ Негативные сценарии",
        "✓ Граничные условия",
        "✓ Валидация данных",
        "✓ Обработка ошибок",
        "✓ Security проверки",
    ]

    for i, item in enumerate(coverage_items):
        add_body_text(slide, item, Inches(8.5), Inches(4.15 + i * 0.4),
                      Inches(3.8), Inches(0.35), font_size=14, font_color=COLORS['dark_gray'])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Мы протестировали систему на PetStore API - стандартном примере REST API.
- Input: 9 endpoints и 2 страницы документации.
- Output: 52 comprehensive тест-кейса за 5 минут.
- Для сравнения: ручное написание такого же количества тест-кейсов заняло бы ~17 часов.
- Это ускорение в 204 раза!
- Покрытие включает не только happy path, но и негативные сценарии, граничные условия, security проверки.
"""


def create_slide_8_pilot_proposal(prs, downtime_data, roi_data):
    """Slide 8: Pilot Proposal"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "4-Week Infrastructure Pilot",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=36, font_color=COLORS['primary_blue'])

    # Timeline (Gantt-style)
    weeks = [
        {"name": "Week 1", "title": "Setup & Training", "color": COLORS['secondary_blue'],
         "tasks": ["Настройка окружения", "Обучение (2 часа)", "Выбор требований"]},
        {"name": "Week 2", "title": "Generation", "color": COLORS['green'],
         "tasks": ["200-300 тест-кейсов", "Ревью и приоритизация", "Адаптация"]},
        {"name": "Week 3", "title": "Execution", "color": COLORS['orange'],
         "tasks": ["Выполнение тестов", "Документирование багов", "Цель: 5-10 багов"]},
        {"name": "Week 4", "title": "Scale Planning", "color": COLORS['primary_blue'],
         "tasks": ["Анализ результатов", "Планирование", "Roadmap 3-6 мес"]},
    ]

    for i, week in enumerate(weeks):
        left = Inches(0.5 + i * 3.2)

        # Week header
        add_rounded_rectangle(slide, left, Inches(1.2), Inches(3), Inches(0.5), week['color'])
        add_body_text(slide, week['name'], left + Inches(0.1), Inches(1.27),
                      Inches(2.8), Inches(0.4), font_size=14, font_color=COLORS['white'], bold=True)
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Week content box
        add_rounded_rectangle(slide, left, Inches(1.75), Inches(3), Inches(2.4), COLORS['light_blue'])

        # Week title
        add_title_shape(slide, week['title'], left + Inches(0.15), Inches(1.85),
                        Inches(2.7), Inches(0.5), font_size=16, font_color=COLORS['primary_blue'])

        # Tasks
        for j, task in enumerate(week['tasks']):
            add_body_text(slide, f"• {task}", left + Inches(0.15), Inches(2.35 + j * 0.5),
                          Inches(2.7), Inches(0.45), font_size=12, font_color=COLORS['dark_gray'])

    # Success metrics
    add_title_shape(slide, "Success Metrics:", Inches(0.5), Inches(4.4),
                    Inches(5), Inches(0.5), font_size=20, font_color=COLORS['primary_blue'])

    metrics = [
        ("200+", "тест-кейсов\nсгенерировано"),
        ("5-10", "реальных багов\nнайдено"),
        ("0", "инцидентов\nво время пилота"),
        ("8+", "удовлетворённость\nинженеров (/10)"),
    ]

    for i, (number, label) in enumerate(metrics):
        left = Inches(0.5 + i * 3.2)

        add_rounded_rectangle(slide, left, Inches(4.9), Inches(3), Inches(1.7), COLORS['light_blue'])

        add_title_shape(slide, number, left + Inches(0.1), Inches(5),
                        Inches(2.8), Inches(0.6), font_size=36, font_color=COLORS['green'],
                        alignment=PP_ALIGN.CENTER)

        add_body_text(slide, label, left + Inches(0.1), Inches(5.55),
                      Inches(2.8), Inches(0.6), font_size=12, font_color=COLORS['gray'])
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Resource requirement
    add_rounded_rectangle(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), COLORS['yellow'])
    add_body_text(slide, "Требуется: 2 инженера (part-time, ~8 часов/неделя) | Budget: 0₽ (MVP готов)",
                  Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
                  font_size=14, font_color=COLORS['dark_gray'], bold=True)
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Предлагаем 4-недельный пилот для Infrastructure направления.
- Week 1: Настройка и обучение - всего 2 часа на обучение инженеров.
- Week 2: Генерация 200-300 тест-кейсов для приоритетных компонентов.
- Week 3: Выполнение тестов, цель - найти 5-10 реальных багов.
- Week 4: Анализ результатов и планирование масштабирования.
- Success metrics чёткие и измеримые.
- Требуется всего 2 инженера part-time. Budget = 0, потому что MVP уже готов.
"""


def create_slide_9_roi(prs, downtime_data, roi_data):
    """Slide 9: ROI & Economics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Return on Investment",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=40, font_color=COLORS['primary_blue'])

    # Option A: Hire QA
    add_rounded_rectangle(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5), COLORS['light_blue'])
    add_rectangle(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.6), COLORS['red'])

    add_title_shape(slide, "Option A: Нанять QA", Inches(0.7), Inches(1.28),
                    Inches(5.4), Inches(0.5), font_size=20, font_color=COLORS['white'])

    hire_data = [
        ("Recruitment:", "6 месяцев"),
        ("Onboarding:", "3 месяца"),
        ("Год 1:", "~5.0M₽"),
        ("Год 2:", "~5.1M₽"),
        ("Год 3:", "~5.7M₽"),
    ]

    for i, (label, value) in enumerate(hire_data):
        add_body_text(slide, label, Inches(0.7), Inches(1.95 + i * 0.5),
                      Inches(2.5), Inches(0.4), font_size=14, font_color=COLORS['gray'])
        add_body_text(slide, value, Inches(3.2), Inches(1.95 + i * 0.5),
                      Inches(2.8), Inches(0.4), font_size=14, font_color=COLORS['dark_gray'], bold=True)

    add_rectangle(slide, Inches(0.7), Inches(4.3), Inches(5.4), Inches(0.05), COLORS['gray'])
    add_title_shape(slide, "ИТОГО 3 года: ~15.8M₽", Inches(0.7), Inches(4.4),
                    Inches(5.4), Inches(0.4), font_size=18, font_color=COLORS['red'])

    # Option B: AI Generator
    add_rounded_rectangle(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(3.5), COLORS['light_blue'])
    add_rectangle(slide, Inches(6.5), Inches(1.2), Inches(6.3), Inches(0.6), COLORS['green'])

    add_title_shape(slide, "Option B: AI Generator", Inches(6.7), Inches(1.28),
                    Inches(6), Inches(0.5), font_size=20, font_color=COLORS['white'])

    ai_data = [
        ("Development:", "Готов! ✓"),
        ("Time to value:", "4 недели"),
        ("Год 1:", "~0.4M₽"),
        ("Год 2:", "~0.4M₽"),
        ("Год 3:", "~0.4M₽"),
    ]

    for i, (label, value) in enumerate(ai_data):
        add_body_text(slide, label, Inches(6.7), Inches(1.95 + i * 0.5),
                      Inches(2.8), Inches(0.4), font_size=14, font_color=COLORS['gray'])
        color = COLORS['green'] if "Готов" in value or "4 недели" in value else COLORS['dark_gray']
        add_body_text(slide, value, Inches(9.5), Inches(1.95 + i * 0.5),
                      Inches(3), Inches(0.4), font_size=14, font_color=color, bold=True)

    add_rectangle(slide, Inches(6.7), Inches(4.3), Inches(5.9), Inches(0.05), COLORS['gray'])
    add_title_shape(slide, "ИТОГО 3 года: ~1.2M₽", Inches(6.7), Inches(4.4),
                    Inches(5.9), Inches(0.4), font_size=18, font_color=COLORS['green'])

    # Big numbers summary
    summary_data = [
        ("14.6M₽", "ЭКОНОМИЯ\nза 3 года", COLORS['green']),
        ("1,200%", "ROI", COLORS['primary_blue']),
        ("9x", "быстрее\nдо результата", COLORS['orange']),
    ]

    for i, (number, label, color) in enumerate(summary_data):
        left = Inches(0.5 + i * 4.4)

        add_rounded_rectangle(slide, left, Inches(5), Inches(4.1), Inches(1.9), color)

        add_title_shape(slide, number, left + Inches(0.1), Inches(5.15),
                        Inches(3.9), Inches(0.8), font_size=40, font_color=COLORS['white'],
                        alignment=PP_ALIGN.CENTER)

        add_body_text(slide, label, left + Inches(0.1), Inches(5.9),
                      Inches(3.9), Inches(0.6), font_size=14, font_color=COLORS['white'], bold=True)
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Давайте посмотрим на экономику двух вариантов.
- Option A (Найм QA): 6 месяцев на поиск + 3 месяца онбординг. За 3 года = ~15.8M рублей.
- Option B (AI Generator): MVP готов сейчас, 4 недели до результата. За 3 года = ~1.2M рублей.
- Экономия за 3 года: 14.6 миллиона рублей.
- ROI: 1,200%.
- И главное: результат в 9 раз быстрее - 4 недели против 9 месяцев.
"""


def create_slide_10_risk_timeline(prs, downtime_data, roi_data):
    """Slide 10: Risk Mitigation Timeline"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Путь к снижению риска",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=40, font_color=COLORS['primary_blue'])

    # Chart
    chart_path = CHARTS_DIR / 'risk_timeline.png'
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.2), width=Inches(12))

    # Detailed timeline boxes
    timeline_items = [
        {
            "period": "Сегодня",
            "status": "CRITICAL",
            "color": COLORS['red'],
            "coverage": "0%",
            "incidents": "Высокий риск",
            "confidence": "Низкая"
        },
        {
            "period": "Месяц 1",
            "status": "MEDIUM",
            "color": COLORS['yellow'],
            "coverage": "20%",
            "incidents": "Контролируемый",
            "confidence": "Растёт"
        },
        {
            "period": "Месяц 3",
            "status": "LOW",
            "color": COLORS['green'],
            "coverage": "60%",
            "incidents": "Минимальный",
            "confidence": "Высокая"
        },
        {
            "period": "Месяц 6",
            "status": "CONTROLLED",
            "color": COLORS['green'],
            "coverage": "80%+",
            "incidents": "Управляемый",
            "confidence": "Полная"
        },
    ]

    for i, item in enumerate(timeline_items):
        left = Inches(0.5 + i * 3.2)
        top = Inches(3.8)

        # Header with status
        add_rectangle(slide, left, top, Inches(3), Inches(0.5), item['color'])
        add_body_text(slide, f"{item['period']}: {item['status']}", left + Inches(0.1), top + Inches(0.1),
                      Inches(2.8), Inches(0.35), font_size=12, font_color=COLORS['white'], bold=True)
        slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content box
        add_rounded_rectangle(slide, left, top + Inches(0.55), Inches(3), Inches(1.8), COLORS['light_blue'])

        # Metrics
        metrics = [
            f"Покрытие: {item['coverage']}",
            f"Риск инцидента: {item['incidents']}",
            f"Уверенность: {item['confidence']}"
        ]

        for j, metric in enumerate(metrics):
            add_body_text(slide, metric, left + Inches(0.15), top + Inches(0.7 + j * 0.5),
                          Inches(2.7), Inches(0.45), font_size=12, font_color=COLORS['dark_gray'])

    # Bottom insight
    add_rounded_rectangle(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6), COLORS['primary_blue'])
    add_body_text(slide, "Каждый месяц без действий = увеличение вероятности инцидента стоимостью 50M₽+",
                  Inches(0.7), Inches(6.42), Inches(12), Inches(0.4),
                  font_size=16, font_color=COLORS['white'], bold=True)
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Сегодня: критический риск, покрытие 0%, любой инцидент может быть катастрофическим.
- Месяц 1 (после пилота): снижаем риск до среднего, покрытие 20% критических компонентов.
- Месяц 3: низкий риск, 60% покрытие, уверенность в стабильности растёт.
- Месяц 6: контролируемый риск, 80%+ покрытие, полная уверенность в качестве.
- Ключевой message: каждый месяц бездействия - это накопление риска потенциального инцидента на 50+ миллионов.
"""


def create_slide_11_cta(prs, downtime_data, roi_data):
    """Slide 11: Call to Action"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['white'])

    # Title
    add_title_shape(slide, "Решение требуется",
                    Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=44, font_color=COLORS['primary_blue'])

    # Approve section (green)
    add_rounded_rectangle(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3.2), COLORS['green'])
    add_title_shape(slide, "✓ Одобрить пилот", Inches(0.7), Inches(1.45),
                    Inches(5.6), Inches(0.5), font_size=24, font_color=COLORS['white'])

    approve_items = [
        "Длительность: 4 недели",
        "Ресурсы: 2 инженера part-time",
        "Budget: 0₽ (MVP готов)",
        "Результат: 200+ тест-кейсов,",
        "   5-10 найденных багов",
    ]

    for i, item in enumerate(approve_items):
        add_body_text(slide, item, Inches(0.9), Inches(2.1 + i * 0.45),
                      Inches(5.2), Inches(0.4), font_size=16, font_color=COLORS['white'])

    # Reject section (red)
    add_rounded_rectangle(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3.2), COLORS['red'])
    add_title_shape(slide, "✗ Альтернатива", Inches(7), Inches(1.45),
                    Inches(5.6), Inches(0.5), font_size=24, font_color=COLORS['white'])

    reject_items = [
        "Продолжать без покрытия",
        "Накапливать технический долг",
        "Надеяться на удачу",
        "Риск инцидента: $5M+/час",
        "Штрафы регулятора: £50M+",
    ]

    for i, item in enumerate(reject_items):
        add_body_text(slide, item, Inches(7.2), Inches(2.1 + i * 0.45),
                      Inches(5.4), Inches(0.4), font_size=16, font_color=COLORS['white'])

    # Next steps
    add_title_shape(slide, "Следующий шаг:", Inches(0.5), Inches(4.8),
                    Inches(5), Inches(0.5), font_size=24, font_color=COLORS['primary_blue'])

    add_rounded_rectangle(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4), COLORS['yellow'])

    next_steps = [
        "1. Kickoff meeting на следующей неделе",
        "2. Выбор первых компонентов для тестирования",
        "3. Старт пилота - генерация первых тест-кейсов",
    ]

    for i, step in enumerate(next_steps):
        add_body_text(slide, step, Inches(0.7), Inches(5.55 + i * 0.4),
                      Inches(11.8), Inches(0.35), font_size=16, font_color=COLORS['dark_gray'], bold=True)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Нам нужно принять решение сегодня.
- Вариант А: одобрить 4-недельный пилот. Ресурсы минимальны - 2 инженера part-time, budget = 0.
- Вариант Б: продолжать как есть. Это означает накопление риска с каждым днём.
- Напоминаю: стоимость часа простоя в финансовом секторе - более $5M. Штрафы регулятора - от £50M.
- Следующий шаг: kickoff meeting на следующей неделе. Готов назначить сразу после этой встречи.
"""


def create_slide_12_contact(prs, downtime_data, roi_data):
    """Slide 12: Contact / Thank You"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['primary_blue'])

    # Main question
    add_title_shape(slide, "Вопросы?",
                    Inches(0.5), Inches(2), Inches(12), Inches(1),
                    font_size=72, font_color=COLORS['white'], bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Contact info box
    add_rounded_rectangle(slide, Inches(3), Inches(3.5), Inches(7), Inches(2.2), COLORS['light_blue'])

    add_title_shape(slide, "AI Test Generator", Inches(3.2), Inches(3.7),
                    Inches(6.6), Inches(0.5), font_size=24, font_color=COLORS['primary_blue'],
                    alignment=PP_ALIGN.CENTER)

    add_body_text(slide, "GitHub: github.com/snromanov/ai-test-generator",
                  Inches(3.2), Inches(4.3), Inches(6.6), Inches(0.4),
                  font_size=16, font_color=COLORS['dark_gray'])
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_body_text(slide, "Tech Stack: Python, LangChain, Claude API",
                  Inches(3.2), Inches(4.7), Inches(6.6), Inches(0.4),
                  font_size=14, font_color=COLORS['gray'])
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_body_text(slide, "On-premise deployment ready",
                  Inches(3.2), Inches(5.1), Inches(6.6), Inches(0.4),
                  font_size=14, font_color=COLORS['gray'])
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer
    add_body_text(slide, "Январь 2025 | Райффайзенбанк",
                  Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                  font_size=14, font_color=COLORS['light_blue'])
    slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
TALKING POINTS:
- Спасибо за внимание!
- Готов ответить на любые вопросы.
- Код доступен на GitHub, можем провести демо прямо сейчас если интересно.
- Ещё раз: MVP готов, можем начать пилот на следующей неделе.
"""


def generate_presentation():
    """Main function to generate the complete presentation."""
    print("Loading data...")
    downtime_data, roi_data = load_data()

    print("Creating charts...")
    create_charts(downtime_data, roi_data)

    print("Generating presentation...")

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Generate all slides
    create_slide_1_title(prs, downtime_data, roi_data)
    create_slide_2_problem(prs, downtime_data, roi_data)
    create_slide_3_downtime_cost(prs, downtime_data, roi_data)
    create_slide_4_real_incidents(prs, downtime_data, roi_data)
    create_slide_5_traditional_solutions(prs, downtime_data, roi_data)
    create_slide_6_solution(prs, downtime_data, roi_data)
    create_slide_7_demo_results(prs, downtime_data, roi_data)
    create_slide_8_pilot_proposal(prs, downtime_data, roi_data)
    create_slide_9_roi(prs, downtime_data, roi_data)
    create_slide_10_risk_timeline(prs, downtime_data, roi_data)
    create_slide_11_cta(prs, downtime_data, roi_data)
    create_slide_12_contact(prs, downtime_data, roi_data)

    # Save presentation
    OUTPUT_DIR.mkdir(exist_ok=True)
    output_path = OUTPUT_DIR / 'ai_test_generator_pitch.pptx'
    prs.save(str(output_path))

    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    return output_path


if __name__ == '__main__':
    generate_presentation()
